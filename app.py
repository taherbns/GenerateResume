from flask import Flask, request, jsonify
import PyPDF2
from docx import Document
from pptx import Presentation
from transformers import BartTokenizer, BartForConditionalGeneration, T5Tokenizer, T5ForConditionalGeneration

app = Flask(__name__)

# Initialiser les modèles BART et T5
bart_model = BartForConditionalGeneration.from_pretrained('facebook/bart-large-cnn')
bart_tokenizer = BartTokenizer.from_pretrained('facebook/bart-large-cnn')

t5_model = T5ForConditionalGeneration.from_pretrained('t5-base')
t5_tokenizer = T5Tokenizer.from_pretrained('t5-base', legacy=False)

# Fonction pour choisir le modèle (BART ou T5)
def select_model(model_name):
    if model_name == 'bart':
        return bart_model, bart_tokenizer
    elif model_name == 't5':
        return t5_model, t5_tokenizer
    else:
        raise ValueError("Modèle non supporté : choisir 'bart' ou 't5'")

# Fonction d'extraction de texte depuis les fichiers
def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, "rb") as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text += page.extract_text()
    return text

def extract_text_from_word(file_path):
    text = ""
    doc = Document(file_path)
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_text_from_ppt(file_path):
    text = ""
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Fonction pour ajuster la longueur maximale et minimale du résumé
def ajuster_longueur(text, proportion=0.3, max_cap=2000, min_cap=20):
    input_length = len(text.split())  # Nombre de mots dans le texte
    max_length = int(input_length * proportion)  # Longueur maximale proportionnelle
    max_length = min(max_length, max_cap)  # Limite maximale à 200 tokens
    min_length = min(min_cap, max_length - 1)  # Assurer que min_length < max_length
    return max_length, min_length

# Fonction pour diviser le texte en morceaux respectant la limite de tokens
def split_text(text, max_tokens=2000):
    words = text.split()
    chunks = []
    current_chunk = []

    for word in words:
        current_chunk.append(word)
        if len(current_chunk) >= max_tokens:
            chunks.append(" ".join(current_chunk))
            current_chunk = []

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    return chunks

# Fonction de génération de résumé avec BART ou T5
def summarize_text(text, model_name='bart'):
    model, tokenizer = select_model(model_name)

    # Diviser le texte en segments s'il est trop long
    max_chunk_size = 1024
    text_chunks = [text[i:i + max_chunk_size] for i in range(0, len(text), max_chunk_size)]
    
    summary = ""
    for chunk in text_chunks:
        max_length, min_length = ajuster_longueur(chunk)
        
        inputs = tokenizer.encode(chunk, return_tensors="pt", max_length=1024, truncation=True)
        summary_ids = model.generate(inputs, max_length=max_length, min_length=min_length, length_penalty=2.0, num_beams=4, early_stopping=True)
        
        chunk_summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
        summary += chunk_summary + " "

    return summary

# Route pour téléverser un fichier et générer un résumé
@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "Aucun fichier n'a été trouvé"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "Aucun fichier sélectionné"}), 400

    file_path = f"uploaded_file.{file.filename.rsplit('.', 1)[1].lower()}"
    file.save(file_path)
    
    print(f"Fichier reçu : {file.filename}")  # Log pour vérifier le fichier reçu
    
    try:
        # Déterminer le type de fichier et extraire le texte
        if file_path.endswith('.pdf'):
            text = extract_text_from_pdf(file_path)
        elif file_path.endswith('.docx'):
            text = extract_text_from_word(file_path)
        elif file_path.endswith('.pptx'):
            text = extract_text_from_ppt(file_path)
        else:
            return jsonify({"error": "Type de fichier non supporté"}), 400

        # Log pour vérifier si le texte a été extrait
        print(f"Texte extrait : {text[:500]}...")  # Afficher les 500 premiers caractères du texte extrait
        
        # Vérifier si le texte extrait est trop court pour être résumé
        if len(text.strip()) == 0:
            return jsonify({"error": "Le fichier est vide ou le texte n'a pas pu être extrait"}), 400
        
        # Résumer le texte avec le modèle choisi
        model_name = request.form.get('model', 'bart')  # par défaut, utiliser BART
        summary = summarize_text(text, model_name=model_name)
        
        print(f"Résumé généré : {summary}")  # Log pour vérifier le résumé généré
        
        return jsonify({"summary": summary})
    
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    app.run(debug=True)
